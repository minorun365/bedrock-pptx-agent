# 必要なPythonライブラリをインポート
import os, json, boto3
from pptx import Presentation
from datetime import datetime

# 環境変数からS3バケット名を取得
S3_BUCKET_NAME = os.getenv("S3_BUCKET_NAME")

# Lambdaのメイン関数
def lambda_handler(event, context):

    # イベントパラメータからトピックとコンテンツを抽出
    topic = next((item["value"] for item in event["parameters"] if item["name"] == "topic"), "")
    content = next((item["value"] for item in event["parameters"] if item["name"] == "content"), "")
    
    # 空白文字を削除し、コンテンツを空行で分割
    content = content.strip()
    slides_content = content.split('\n\n')
    
    # プレゼンテーションオブジェクトを作成
    prs = Presentation()
    
    # タイトルスライドの作成
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"
    
    # コンテンツスライドの作成
    for i, slide_content in enumerate(slides_content):
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        # 見出しと本文を作成
        lines = slide_content.split('\n')
        title.text = lines[0].lstrip('- ')
        content_shape.text = '\n'.join([line.lstrip('- ') for line in lines[1:]])
    
    # S3にファイルを保存する準備
    s3 = boto3.client("s3")
    bucket_name = S3_BUCKET_NAME
    file_name = f"{topic.replace(' ', '_')}.pptx"
    file_path = f"/tmp/{file_name}"

    # S3バケットにファイルをアップロード
    prs.save(file_path)
    s3.upload_file(file_path, bucket_name, file_name)
    
    # ファイルへの署名付きURLを生成
    url = s3.generate_presigned_url(
        'get_object',
        Params={'Bucket': bucket_name, 'Key': file_name},
        ExpiresIn=3600
    )
    
    # エージェント用のレスポンスを返す
    return {
        "messageVersion": "1.0",
        "response": {
            "actionGroup": event["actionGroup"],
            "function": event["function"],
            "functionResponse": {
                "responseBody": {
                    "TEXT": {
                        "body": json.dumps(
                            {"signed_url": url}
                        )
                    }
                }
            },
        },
    }